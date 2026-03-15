"""
document_processor.py  ─  Week 3-4 Enhanced
Multi-format processor with:
  • Chapter / section heading detection
  • Chunks returned as dicts {"text": ..., "chapter": ...}
  • Subject hinted from filename + content
"""

import os
import re
from typing import Optional


# ─── Content-type keyword heuristics ─────────────────────────────────────────
_TYPE_KEYWORDS = {
    "question_paper": [
        r"\bq\d+\b", r"\bquestion\s+\d+\b", r"\bmark[s]?\b",
        r"\bexam\b", r"\buniversity\b", r"\btime\s*:\s*\d",
        r"\banswer\s+all\b", r"\bsection\s+[a-z]\b",
    ],
    "lab_manual": [
        r"\baim\b", r"\bapparatus\b", r"\bprocedure\b",
        r"\bobservation\b", r"\bviva\b", r"\bexperiment\b",
        r"\bresult\b.*\btable\b",
    ],
    "textbook": [
        r"\bchapter\s+\d+\b", r"\bdefinition\b", r"\btheorem\b",
        r"\bexercise\b", r"\bsummary\b", r"\blearning\s+objective",
    ],
    "notes": [
        r"\blecture\s+\d+\b", r"\bnote[s]?\b", r"\btopic\s*:\s*",
        r"\bunit\s+\d+\b", r"\bimportant\b",
    ],
}

# ─── Heading detection patterns ───────────────────────────────────────────────
_HEADING_PATTERNS = [
    # "Chapter 1: Title"  /  "Unit 2 – Title"  /  "Module 3 - Title"
    re.compile(
        r"^(?:chapter|unit|module|lecture|section|topic)\s+[\dIVXivx]+[\s:\-–]+(.{3,80})$",
        re.IGNORECASE | re.MULTILINE,
    ),
    # "1.  Title"  /  "1.2  Sub-title"  (capitalised, ≥4 chars)
    re.compile(
        r"^\s*\d+(?:\.\d+)*[\s\.\)]{1,3}([A-Z][A-Za-z\s\-]{3,60})\s*$",
        re.MULTILINE,
    ),
]


def _detect_content_type(text: str) -> str:
    text_lower = text.lower()
    scores = {t: 0 for t in _TYPE_KEYWORDS}
    for ctype, patterns in _TYPE_KEYWORDS.items():
        for pat in patterns:
            scores[ctype] += len(re.findall(pat, text_lower))
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "general"


def _extract_headings(text: str) -> list[tuple[int, str]]:
    """Return sorted list of (char_position, heading_title)."""
    found = []
    for pat in _HEADING_PATTERNS:
        for m in pat.finditer(text):
            title = (m.group(1) if m.lastindex else m.group(0)).strip()
            title = re.sub(r"\s+", " ", title)[:80]
            found.append((m.start(), title))
    # Sort and deduplicate by position
    found.sort(key=lambda x: x[0])
    seen, unique = set(), []
    for pos, title in found:
        if pos not in seen:
            seen.add(pos)
            unique.append((pos, title))
    return unique


def _chapter_at(pos: int, headings: list[tuple[int, str]]) -> str:
    chapter = "General"
    for h_pos, h_title in headings:
        if h_pos <= pos:
            chapter = h_title
        else:
            break
    return chapter


def _chunk_text(
    text: str, chunk_size: int = 800, overlap: int = 150
) -> list[dict]:
    """
    Split text into overlapping chunks, each tagged with its chapter/section.
    Returns list of {"text": str, "chapter": str}.
    """
    headings   = _extract_headings(text)
    sentences  = re.split(r"(?<=[.!?])\s+", text)
    chunks: list[dict] = []
    current: list[str] = []
    current_len = 0
    char_pos    = 0

    for sent in sentences:
        slen = len(sent)
        if current_len + slen > chunk_size and current:
            body = " ".join(current).strip()
            if body:
                chap = _chapter_at(char_pos - current_len, headings)
                chunks.append({"text": body, "chapter": chap})
            # overlap carry-forward
            overlap_chars, overlap_sents = 0, []
            for s in reversed(current):
                overlap_chars += len(s)
                overlap_sents.insert(0, s)
                if overlap_chars >= overlap:
                    break
            current     = overlap_sents
            current_len = sum(len(s) for s in current)

        current.append(sent)
        current_len += slen
        char_pos    += slen + 1

    if current:
        body = " ".join(current).strip()
        if body:
            chap = _chapter_at(char_pos - current_len, headings)
            chunks.append({"text": body, "chapter": chap})

    return chunks


# ─── Format extractors ────────────────────────────────────────────────────────
def _extract_pdf(filepath: str) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        return "\n".join(parts)
    except ImportError:
        pass
    try:
        import PyPDF2
        parts = []
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def _extract_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc   = Document(filepath)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def _extract_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(filepath)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def _extract_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ─── Main ─────────────────────────────────────────────────────────────────────
def process_document(filepath: str, force_type: Optional[str] = None) -> dict:
    """
    Process a document and return:
        filename, content_type, raw_text,
        chunks (list of {"text", "chapter"}),
        num_chunks, file_type
    """
    filename = os.path.basename(filepath)
    ext      = os.path.splitext(filename)[1].lower()

    extractors = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".txt":  _extract_txt,
        ".md":   _extract_txt,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported file type: {ext}")

    raw_text     = extractors[ext](filepath)
    content_type = force_type or _detect_content_type(raw_text)
    chunks       = _chunk_text(raw_text)

    return {
        "filename":     filename,
        "content_type": content_type,
        "raw_text":     raw_text,
        "chunks":       chunks,
        "num_chunks":   len(chunks),
        "file_type":    ext,
    }